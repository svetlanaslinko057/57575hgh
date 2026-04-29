"""
File → text extractor for the /estimate flow.

Accepts popular office/doc/image formats from web + mobile clients and
returns a single UTF-8 text blob the estimator can treat like a pasted brief.

Supported:
  • Plain text / markdown (.txt, .md)
  • PDF                    — pypdf
  • DOCX (Open XML)        — python-docx
  • XLSX (Open XML)        — openpyxl
  • PPTX (Open XML)        — python-pptx
  • Images                 — pytesseract (eng + rus) on .png/.jpg/.jpeg/.webp
  • HEIC                   — Pillow-HEIF if available, otherwise 415

Legacy binary MS formats (.doc/.xls/.ppt) are refused with a friendly
message — we won't ship libreoffice just for a brief.

Limits:
  • 10 MB max upload (enough for a brief PDF / slide deck)
  • 8000 chars max extracted text (more than enough for an estimate)
"""
from __future__ import annotations

import io
import logging
import re
from typing import Tuple

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel

logger = logging.getLogger("file_parser")

MAX_BYTES = 10 * 1024 * 1024           # 10 MB
MAX_TEXT_CHARS = 8000                  # keep the estimate prompt sane


class ParsedFile(BaseModel):
    name: str
    size: int
    mime: str
    text: str
    truncated: bool = False
    source: str  # "text", "pdf", "docx", "xlsx", "pptx", "image-ocr"


# ── extractors ───────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # Tables
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
    return "\n".join(parts)


def _extract_image_ocr(data: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
    except Exception as e:
        raise HTTPException(
            status_code=415,
            detail="Image OCR is not available on this server. Try a text file or PDF.",
        ) from e
    try:
        img = Image.open(io.BytesIO(data))
        # Tesseract handles most RGB/gray modes; convert everything else.
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        # Dual-lang OCR — English + Russian cover most briefs.
        text = pytesseract.image_to_string(img, lang="eng+rus")
        return text
    except pytesseract.TesseractNotFoundError as e:
        raise HTTPException(
            status_code=415,
            detail="Image OCR engine not installed.",
        ) from e
    except Exception as e:
        logger.warning(f"OCR failed: {e}")
        raise HTTPException(
            status_code=422,
            detail=f"Could not read text from this image: {e}",
        ) from e


# ── MIME / extension routing ─────────────────────────────────────────────────

def _dispatch(filename: str, mime: str, data: bytes) -> Tuple[str, str]:
    """Return (extracted_text, source_label)."""
    name = (filename or "").lower()
    mime = (mime or "").lower()

    # Text first — cheapest.
    if mime.startswith("text/") or name.endswith((".txt", ".md", ".markdown")):
        try:
            return data.decode("utf-8", errors="replace"), "text"
        except Exception:
            return data.decode("latin-1", errors="replace"), "text"

    if name.endswith(".pdf") or mime == "application/pdf":
        return _extract_pdf(data), "pdf"

    if name.endswith(".docx") or mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return _extract_docx(data), "docx"

    if name.endswith(".xlsx") or mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ):
        return _extract_xlsx(data), "xlsx"

    if name.endswith(".pptx") or mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return _extract_pptx(data), "pptx"

    if mime.startswith("image/") or name.endswith(
        (".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif", ".bmp", ".gif")
    ):
        return _extract_image_ocr(data), "image-ocr"

    # Legacy binary MS formats — we don't support them.
    if name.endswith((".doc", ".xls", ".ppt")):
        raise HTTPException(
            status_code=415,
            detail=(
                f"Legacy {name.rsplit('.', 1)[-1].upper()} format is not "
                "supported yet. Save as .docx / .xlsx / .pptx or PDF and try again."
            ),
        )

    raise HTTPException(
        status_code=415,
        detail=f"Unsupported file type '{mime or name}'. "
               "Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, PNG/JPG/WEBP images.",
    )


# ── router ───────────────────────────────────────────────────────────────────

router = APIRouter()


@router.post("/estimate/parse-file", response_model=ParsedFile)
async def parse_uploaded_brief(file: UploadFile = File(...)):
    """Read a file from multipart upload and return its extracted text.

    Auth-free on purpose — same as `/estimate`, this powers the pre-login
    "we already calculated your product" experience. Abuse is bounded by the
    10 MB cap.
    """
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file")
    if len(data) > MAX_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large ({len(data) // (1024*1024)} MB). Max 10 MB.",
        )

    text, source = _dispatch(
        filename=file.filename or "",
        mime=(file.content_type or "").split(";")[0].strip(),
        data=data,
    )

    # Normalize whitespace — collapse runs of blank lines, strip NULs, rstrip.
    text = text.replace("\x00", "").strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+\n", "\n", text)

    truncated = False
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS].rstrip() + " …[truncated]"
        truncated = True

    if not text or not text.strip():
        raise HTTPException(
            status_code=422,
            detail=(
                "Couldn't extract readable text from this file. "
                "If it's a scanned image, try a clearer version."
            ),
        )

    logger.info(
        f"PARSE-FILE: {file.filename!r} mime={file.content_type} "
        f"bytes={len(data)} source={source} chars={len(text)} "
        f"truncated={truncated}"
    )
    return ParsedFile(
        name=file.filename or "file",
        size=len(data),
        mime=(file.content_type or "application/octet-stream"),
        text=text,
        truncated=truncated,
        source=source,
    )
